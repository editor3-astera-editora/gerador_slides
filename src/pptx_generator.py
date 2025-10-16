"""
Módulo para geração de apresentações PowerPoint a partir de slides JSON.
Suporta renderização de fórmulas matemáticas em LaTeX.
"""

import json
import re
from pathlib import Path
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN
import matplotlib.pyplot as plt
import matplotlib
matplotlib.use('Agg')  # Backend sem GUI
import io


def detectar_latex(texto):
    """
    Detecta se um texto contém LaTeX ($ ou $$).

    Args:
        texto (str): Texto a verificar

    Returns:
        bool: True se contém LaTeX
    """
    return '$' in texto


def separar_texto_e_latex(bullet):
    """
    Separa um bullet em parte de texto e parte LaTeX.

    Exemplo:
    "Fórmula do vértice: $t = -b/(2a)$" → ("Fórmula do vértice: ", "$t = -b/(2a)$")

    Args:
        bullet (str): Texto do bullet

    Returns:
        tuple: (texto_antes, formula_latex) ou (bullet, None) se não houver LaTeX
    """
    import re

    # Busca por padrão: texto antes de $...$ ou $$...$$
    match = re.search(r'(.*?)(\$\$.*?\$\$|\$.*?\$)', bullet)

    if match:
        texto_antes = match.group(1).strip()
        formula = match.group(2).strip()
        return (texto_antes, formula)
    else:
        return (bullet, None)


def render_latex_as_image(latex_str, filepath, images_folder):
    """
    Renderiza uma string LaTeX para um arquivo de imagem PNG transparente.

    Args:
        latex_str (str): String LaTeX (com ou sem delimitadores $)
        filepath (Path): Caminho para salvar a imagem
        images_folder (Path): Pasta onde salvar imagens de fórmulas

    Returns:
        io.BytesIO: Buffer da imagem PNG
    """
    try:
        # Cria pasta de imagens se não existir
        images_folder.mkdir(parents=True, exist_ok=True)

        # Remove os delimitadores '$' para processar
        latex_limpo = latex_str.strip()
        # Remove $ externo, mantém conteúdo
        if latex_limpo.startswith('$$') and latex_limpo.endswith('$$'):
            latex_limpo = latex_limpo[2:-2].strip()
        elif latex_limpo.startswith('$') and latex_limpo.endswith('$'):
            latex_limpo = latex_limpo[1:-1].strip()

        print(f"    [DEBUG] Renderizando LaTeX: {latex_limpo}")

        # Configuração do matplotlib para renderização matemática
        plt.rcParams.update({
            'text.usetex': False,
            'mathtext.fontset': 'cm',
            'mathtext.default': 'regular',
            'font.size': 14
        })

        # Cria figura
        fig, ax = plt.subplots(figsize=(12, 2), dpi=150)
        ax.axis('off')

        # Renderiza LaTeX - o matplotlib usa $ para delimitar math mode
        # Adiciona os $ de volta para o matplotlib processar
        ax.text(0.5, 0.5,
               f'${latex_limpo}$',
               fontsize=36,
               ha='center',
               va='center',
               transform=ax.transAxes)

        # Remove margens
        fig.patch.set_alpha(0)  # Transparente
        plt.tight_layout(pad=0.5)

        # Salva em buffer
        img_buffer = io.BytesIO()
        plt.savefig(img_buffer,
                   format='png',
                   transparent=True,
                   bbox_inches='tight',
                   pad_inches=0.3,
                   dpi=150)
        plt.close(fig)
        img_buffer.seek(0)

        # Salva também em arquivo para debug
        image_path = images_folder / f"{filepath.stem}.png"
        with open(image_path, 'wb') as f:
            f.write(img_buffer.getvalue())
        print(f"    [DEBUG] Imagem salva: {image_path}")
        img_buffer.seek(0)

        return img_buffer

    except Exception as e:
        print(f"    ERRO ao renderizar LaTeX '{latex_str}': {e}")
        import traceback
        traceback.print_exc()
        raise


def carregar_slides_json(json_path):
    """
    Carrega os slides do arquivo JSON.

    Args:
        json_path (Path): Caminho para o arquivo JSON de slides

    Returns:
        dict: Dados dos slides
    """
    with open(json_path, 'r', encoding='utf-8') as f:
        return json.load(f)


def criar_template_padrao():
    """
    Cria um template PowerPoint padrão se não existir.
    Template simples com layout título e conteúdo.

    Returns:
        Presentation: Objeto de apresentação padrão
    """
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    return prs


def adicionar_slide_abertura(prs, slide_data):
    """
    Adiciona slide de ABERTURA usando Layout 0 (SOMENTE para o primeiro slide).
    Layout 0: Usa placeholders por idx
    - idx=0: Título
    - idx=1: Texto/conteúdo

    Args:
        prs (Presentation): Objeto de apresentação PowerPoint
        slide_data (dict): Dados do slide (slide_bullets, frase_introdutoria)

    Returns:
        Slide: Slide adicionado
    """
    slide = prs.slides.add_slide(prs.slide_layouts[0])

    titulo = slide_data.get("slide_title", "")
    frase_intro = slide_data.get("frase_introdutoria", "")
    bullets = slide_data.get("slide_bullets", [])

    try:
        # Placeholder idx=0: Título
        if titulo:
            ph_titulo = None
            for ph in slide.placeholders:
                if ph.placeholder_format.idx == 0:
                    ph_titulo = ph
                    break
            if ph_titulo and ph_titulo.has_text_frame:
                # Capitaliza apenas a primeira letra
                titulo_formatado = titulo[0].upper() + titulo[1:].lower() if titulo else ""
                ph_titulo.text = titulo_formatado

        # Placeholder idx=1: Conteúdo (frase introdutória + bullets)
        ph_conteudo = None
        for ph in slide.placeholders:
            if ph.placeholder_format.idx == 1:
                ph_conteudo = ph
                break

        if ph_conteudo and ph_conteudo.has_text_frame:
            text_frame = ph_conteudo.text_frame
            text_frame.clear()

            # Adiciona frase introdutória se existir
            if frase_intro:
                p = text_frame.paragraphs[0]
                p.text = frase_intro
                p.level = 0
                p.font.size = Pt(16)
                p.font.bold = True

            # Adiciona bullets
            for i, bullet in enumerate(bullets):
                if i == 0 and not frase_intro:
                    p = text_frame.paragraphs[0]
                else:
                    p = text_frame.add_paragraph()
                # Adiciona hífen antes do bullet
                p.text = f"- {bullet}"
                p.level = 1 if frase_intro else 0
                p.font.size = Pt(24)

    except Exception as e:
        print(f"    ERRO ao preencher placeholders: {e}")
        print(f"    Placeholders disponíveis no Layout 0:")
        for i, ph in enumerate(slide.placeholders):
            print(f"      [{i}] idx={ph.placeholder_format.idx}, tipo={ph.placeholder_format.type}, nome={ph.name}")
        raise

    return slide


def adicionar_slide_conteudo(prs, slide_data, output_dir=None):
    """
    Adiciona slide de CONTEÚDO usando Layout 1 (usado para TODOS os slides exceto o primeiro).
    Suporta renderização de fórmulas matemáticas em LaTeX.

    Layout 1 atualizado com 3 placeholders:
    - idx=10: Título do conceito (titulo_conceito)
    - idx=1: Frase introdutória
    - idx=2: Bullets (texto ou imagens LaTeX)

    Args:
        prs (Presentation): Objeto de apresentação PowerPoint
        slide_data (dict): Dados do slide (titulo_conceito, frase_introdutoria, slide_bullets)
        output_dir (Path): Diretório de output para salvar imagens de fórmulas

    Returns:
        Slide: Slide adicionado
    """
    slide = prs.slides.add_slide(prs.slide_layouts[1])

    titulo_conceito = slide_data.get("titulo_conceito", "")
    frase_intro = slide_data.get("frase_introdutoria", "")
    bullets = slide_data.get("slide_bullets", [])

    # Layout 1 do template tem:
    # - Placeholder com idx=10: Título do conceito
    # - Placeholder com idx=1: Conteúdo principal (frase introdutória)
    # - Placeholder com idx=2: Texto/bullets

    try:
        # Placeholder idx=10: Título do conceito
        ph_titulo = None
        for ph in slide.placeholders:
            if ph.placeholder_format.idx == 10:
                ph_titulo = ph
                break

        if ph_titulo and ph_titulo.has_text_frame:
            if titulo_conceito:
                # Capitaliza apenas a primeira letra
                titulo_formatado = titulo_conceito[0].upper() + titulo_conceito[1:].lower() if titulo_conceito else ""
                ph_titulo.text = titulo_formatado
            else:
                # Preenche com espaço para evitar placeholder vazio
                ph_titulo.text = " "

        # Placeholder idx=1: Conteúdo principal (frase introdutória)
        ph_conteudo = None
        for ph in slide.placeholders:
            if ph.placeholder_format.idx == 1:
                ph_conteudo = ph
                break

        if ph_conteudo and ph_conteudo.has_text_frame:
            if frase_intro:
                ph_conteudo.text = frase_intro
            else:
                # IMPORTANTE: Preenche com espaço para evitar placeholder vazio
                # Placeholders vazios causam bugs no PowerPoint COM na exportação
                ph_conteudo.text = " "

        # Placeholder idx=2: Bullets (texto ou LaTeX)
        ph_bullets = None
        for ph in slide.placeholders:
            if ph.placeholder_format.idx == 2:
                ph_bullets = ph
                break

        if ph_bullets and ph_bullets.has_text_frame:
            # Verifica se há LaTeX nos bullets
            tem_latex = any(detectar_latex(bullet) for bullet in bullets)

            if tem_latex:
                # Modo LaTeX: usa posicionamento manual com imagens
                # Remove o placeholder de texto
                sp = ph_bullets._element
                sp.getparent().remove(sp)

                # Define pasta de imagens
                if output_dir:
                    images_folder = Path(output_dir) / "imagens_formulas"
                else:
                    images_folder = Path("output") / "imagens_formulas"
                images_folder.mkdir(parents=True, exist_ok=True)

                # Posição inicial (onde estava o placeholder)
                top = ph_bullets.top + Inches(0.1)
                left = ph_bullets.left
                bullet_counter = 0

                for bullet in bullets:
                    if detectar_latex(bullet):
                        # Separa texto de fórmula
                        texto_antes, formula = separar_texto_e_latex(bullet)

                        # Adiciona texto antes (se houver)
                        if texto_antes:
                            txBox = slide.shapes.add_textbox(left, top, Inches(8), Inches(0.4))
                            p = txBox.text_frame.paragraphs[0]
                            p.text = f"- {texto_antes}"
                            p.font.size = Pt(24)
                            top += Inches(0.45)

                        # Renderiza fórmula LaTeX como imagem
                        if formula:
                            try:
                                filepath = Path(f"formula_{bullet_counter}")
                                img_buffer = render_latex_as_image(formula, filepath, images_folder)
                                # Posiciona imagem um pouco à direita (indentação)
                                pic = slide.shapes.add_picture(img_buffer, left + Inches(0.3), top, width=Inches(5.5))
                                top += pic.height + Inches(0.2)
                                bullet_counter += 1
                            except Exception as e:
                                # Fallback: escreve LaTeX como texto puro
                                print(f"    AVISO: Falha ao renderizar LaTeX '{formula}'. Usando texto puro.")
                                txBox = slide.shapes.add_textbox(left + Inches(0.3), top, Inches(7.5), Inches(0.4))
                                p = txBox.text_frame.paragraphs[0]
                                p.text = formula
                                p.font.size = Pt(20)
                                top += Inches(0.5)
                    else:
                        # É texto normal (sem LaTeX)
                        txBox = slide.shapes.add_textbox(left, top, Inches(8), Inches(0.4))
                        p = txBox.text_frame.paragraphs[0]
                        p.text = f"- {bullet}"
                        p.font.size = Pt(24)
                        top += Inches(0.5)
            else:
                # Modo normal (sem LaTeX)
                text_frame = ph_bullets.text_frame
                text_frame.clear()

                for i, bullet in enumerate(bullets):
                    if i == 0:
                        p = text_frame.paragraphs[0]
                    else:
                        p = text_frame.add_paragraph()
                    # Adiciona hífen antes do bullet
                    p.text = f"- {bullet}"
                    p.level = 0
                    p.font.size = Pt(24)

    except Exception as e:
        print(f"    ERRO ao preencher placeholders: {e}")
        print(f"    Placeholders disponíveis no Layout 1:")
        for i, ph in enumerate(slide.placeholders):
            print(f"      [{i}] idx={ph.placeholder_format.idx}, tipo={ph.placeholder_format.type}, nome={ph.name}")
        raise

    return slide


def gerar_apresentacao_powerpoint(slides_json_path, template_path=None, output_dir=None):
    """
    Gera uma apresentação PowerPoint a partir do arquivo JSON de slides.

    Args:
        slides_json_path (Path): Caminho para o arquivo JSON de slides
        template_path (Path): Caminho para template PowerPoint (opcional)
        output_dir (Path, optional): Diretório de output customizado. Padrão: "output/"

    Returns:
        Path: Caminho do arquivo PowerPoint gerado
    """
    print("\n[FASE IV] Geração de PowerPoint")
    print("-" * 70)

    # Carregar dados dos slides
    dados = carregar_slides_json(slides_json_path)
    slides = dados.get("slides", [])

    if not slides:
        print("Aviso: Nenhum slide encontrado no arquivo JSON.")
        return None

    print(f"Carregados {len(slides)} slides do JSON")

    # Criar ou carregar apresentação
    if template_path and template_path.exists():
        print(f"Usando template: {template_path}")
        prs = Presentation(template_path)

        # Remove TODOS os slides do template (começamos do zero)
        # Mantemos apenas os layouts, não os slides de exemplo
        total_slides_template = len(prs.slides)

        if total_slides_template > 0:
            print(f"  Template tem {total_slides_template} slide(s) - removendo todos para começar limpo...")

            # Remove todos os slides do fim para o início
            for idx in reversed(range(len(prs.slides))):
                rId = prs.slides._sldIdLst[idx].rId
                prs.part.drop_rel(rId)
                del prs.slides._sldIdLst[idx]

            print(f"  OK - Template limpo, iniciando com 0 slides")
    else:
        print("Criando apresentação com template padrão")
        prs = criar_template_padrao()

    print(f"Gerando {len(slides)} slides...")

    # Adicionar slides à apresentação
    for i, slide_data in enumerate(slides):
        tipo = slide_data.get("tipo", "")
        titulo = slide_data.get("slide_title", "Sem Título")
        timestamp = slide_data.get("timestamp_inicio", 0.0)

        # LÓGICA SIMPLIFICADA (2 layouts apenas):
        # - Primeiro slide (i==0) → Layout 0 (abertura com título)
        # - Todos os demais slides → Layout 1 (conteúdo sem título)

        if i == 0:
            # Layout 0: SOMENTE para o primeiro slide (abertura)
            print(f"  Slide {i+1}/{len(slides)}: ABERTURA [Layout 0] - {titulo[:40]}... ({timestamp:.2f}s)")
            adicionar_slide_abertura(prs, slide_data)
        else:
            # Layout 1: Todos os demais slides
            print(f"  Slide {i+1}/{len(slides)}: CONTEUDO [Layout 1] - {titulo[:40]}... ({timestamp:.2f}s)")
            adicionar_slide_conteudo(prs, slide_data, output_dir=output_dir)

    # Define diretório de output
    if output_dir is None:
        output_dir = Path("output")
    else:
        output_dir = Path(output_dir)

    # Cria pasta se não existir
    output_dir.mkdir(parents=True, exist_ok=True)

    # Salvar apresentação
    nome_base = slides_json_path.stem.replace("_slides", "")
    output_path = output_dir / f"{nome_base}.pptx"

    prs.save(str(output_path))

    print(f"\nOK Apresentação PowerPoint salva em: {output_path}")
    print(f"OK Total de slides criados: {len(slides)}")

    return output_path


def exibir_info_template(template_path):
    """
    Exibe informações sobre os layouts disponíveis em um template.
    Útil para debug e configuração.

    Args:
        template_path (Path): Caminho para o template PowerPoint
    """
    if not template_path.exists():
        print(f"Template não encontrado: {template_path}")
        return

    prs = Presentation(template_path)

    print(f"\nInformações do Template: {template_path.name}")
    print("=" * 70)
    print(f"Dimensões: {prs.slide_width / Inches(1):.2f}\" x {prs.slide_height / Inches(1):.2f}\"")
    print(f"\nLayouts disponíveis ({len(prs.slide_layouts)}):")

    for i, layout in enumerate(prs.slide_layouts):
        print(f"  [{i}] {layout.name}")

    print("=" * 70)
